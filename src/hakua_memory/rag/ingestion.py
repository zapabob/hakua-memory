"""Document ingestion for PDF, DOCX, PPTX, Markdown, and plain text."""

from __future__ import annotations

import hashlib
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional

from .chunking import chunk_pages, chunk_slides, chunk_text
from .models import Document


def _sha256(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _detect_type(path: Path) -> str:
    suffix = path.suffix.lower()
    return {
        ".pdf": "pdf",
        ".docx": "docx",
        ".pptx": "pptx",
        ".md": "markdown",
        ".markdown": "markdown",
        ".txt": "text",
        ".html": "html",
        ".htm": "html",
    }.get(suffix, "unknown")


def _read_pdf(path: Path) -> tuple[list[str], int]:
    """Read PDF into pages. Requires pypdf or PyMuPDF."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        pages = [page.get_text() for page in doc]
        doc.close()
        return pages, len(pages)
    except ImportError:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return pages, len(pages)
    except ImportError:
        raise ImportError(
            "PDF ingestion requires pypdf or PyMuPDF. Install with: pip install pypdf"
        )


def _read_docx(path: Path) -> tuple[list[str], int]:
    """Read DOCX into paragraphs grouped by page breaks."""
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(str(path))
        paragraphs: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        # DOCX doesn't have explicit page breaks in python-docx reliably,
        # so treat each paragraph as a "page" for chunking purposes.
        # Better: chunk the full text.
        full_text = "\n\n".join(paragraphs)
        return [full_text], 1
    except ImportError:
        raise ImportError(
            "DOCX ingestion requires python-docx. Install with: pip install python-docx"
        )


def _read_pptx(path: Path) -> tuple[list[str], int]:
    """Read PPTX into slides."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides: list[str] = []
        for slide in prs.slides:
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                slides.append("\n".join(slide_texts))
        return slides, len(slides)
    except ImportError:
        raise ImportError(
            "PPTX ingestion requires python-pptx. Install with: pip install python-pptx"
        )


def _read_markdown(path: Path) -> tuple[list[str], int]:
    """Read Markdown, split by headers into sections."""
    text = path.read_text(encoding="utf-8")
    sections = re.split(r"^#{1,6}\s+", text, flags=re.MULTILINE)
    sections = [s.strip() for s in sections if s.strip()]
    return sections, len(sections) if sections else [text]


def _read_text(path: Path) -> tuple[list[str], int]:
    """Read plain text."""
    text = path.read_text(encoding="utf-8")
    return [text], 1


def _read_html(path: Path) -> tuple[list[str], int]:
    """Read HTML, extract text."""
    try:
        from bs4 import BeautifulSoup
        text = path.read_text(encoding="utf-8")
        soup = BeautifulSoup(text, "html.parser")
        return [soup.get_text(separator="\n")], 1
    except ImportError:
        raise ImportError(
            "HTML ingestion requires beautifulsoup4. Install with: pip install beautifulsoup4"
        )


def ingest_document(
    path: Path,
    *,
    title: str = "",
    version: str = "1",
    author: str = "",
    department: str = "",
    created_at: str = "",
    metadata: Optional[dict[str, Any]] = None,
) -> tuple[Document, list]:
    """Ingest a file and return (Document, chunks)."""
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    doc_type = _detect_type(path)
    if doc_type == "unknown":
        raise ValueError(f"Unsupported file type: {path.suffix}")

    raw = path.read_bytes() if doc_type != "markdown" and doc_type != "text" and doc_type != "html" else b""
    content_hash = hashlib.sha256(raw).hexdigest()

    reader = {
        "pdf": _read_pdf,
        "docx": _read_docx,
        "pptx": _read_pptx,
        "markdown": _read_markdown,
        "text": _read_text,
        "html": _read_html,
    }[doc_type]

    sections, count = reader(path)

    doc = Document(
        document_id=f"doc-{uuid.uuid4().hex[:12]}",
        title=title or path.stem,
        source_uri=str(path.resolve()),
        document_type=doc_type,
        version=version,
        author=author,
        department=department,
        created_at=created_at,
        ingested_at=_now_iso(),
        content_hash=content_hash,
        page_count=count if doc_type in ("pdf", "pptx") else 0,
        slide_count=count if doc_type == "pptx" else 0,
        metadata=metadata or {},
    )

    # Chunk based on document type
    if doc_type == "pdf":
        chunks = chunk_pages(sections, doc.document_id)
    elif doc_type == "pptx":
        chunks = chunk_slides(sections, doc.document_id)
    elif doc_type == "markdown":
        chunks = []
        for i, section in enumerate(sections):
            section_chunks = chunk_text(
                section, doc.document_id, section=f"section-{i+1}"
            )
            chunks.extend(section_chunks)
    else:
        chunks = []
        for section in sections:
            chunks.extend(chunk_text(section, doc.document_id))

    return doc, chunks


def ingest_markdown_string(
    text: str,
    *,
    title: str,
    source_uri: str = "",
    version: str = "1",
    author: str = "",
    department: str = "",
    metadata: Optional[dict[str, Any]] = None,
) -> tuple[Document, list]:
    """Ingest a markdown string directly."""
    content_hash = _sha256(text)
    doc = Document(
        document_id=f"doc-{uuid.uuid4().hex[:12]}",
        title=title,
        source_uri=source_uri or f"inline://{content_hash[:16]}",
        document_type="markdown",
        version=version,
        author=author,
        department=department,
        ingested_at=_now_iso(),
        content_hash=content_hash,
        metadata=metadata or {},
    )
    sections = re.split(r"^#{1,6}\s+", text, flags=re.MULTILINE)
    sections = [s.strip() for s in sections if s.strip()]
    chunks = []
    for i, section in enumerate(sections):
        section_chunks = chunk_text(
            section, doc.document_id, section=f"section-{i+1}"
        )
        chunks.extend(section_chunks)
    return doc, chunks


def ingest_text_string(
    text: str,
    *,
    title: str,
    source_uri: str = "",
    document_type: str = "text",
    version: str = "1",
    author: str = "",
    department: str = "",
    metadata: Optional[dict[str, Any]] = None,
) -> tuple[Document, list]:
    """Ingest a plain text string."""
    content_hash = _sha256(text)
    doc = Document(
        document_id=f"doc-{uuid.uuid4().hex[:12]}",
        title=title,
        source_uri=source_uri or f"inline://{content_hash[:16]}",
        document_type=document_type,
        version=version,
        author=author,
        department=department,
        ingested_at=_now_iso(),
        content_hash=content_hash,
        metadata=metadata or {},
    )
    chunks = chunk_text(text, doc.document_id)
    return doc, chunks
